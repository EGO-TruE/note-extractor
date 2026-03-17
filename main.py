"""
智能笔记提取器
用法：python main.py <文件路径或文件夹路径>
"""

import os
import sys
import json
import re
import time
import copy
import csv as csv_module
import subprocess
import tempfile
from pathlib import Path

import requests
from dotenv import load_dotenv

# ─────────────────────────────────────────────
# Section 1: 配置加载
# ─────────────────────────────────────────────

load_dotenv()

API_BASE_URL = os.getenv("API_BASE_URL", "https://www.traxnode.com/v1")
API_KEY      = os.getenv("API_KEY")
MODEL        = os.getenv("MODEL", "gpt-4o")
OUTPUT_DIR   = Path("output")


# ─────────────────────────────────────────────
# Section 2: 文件读取
# ─────────────────────────────────────────────

def read_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- 第 {i + 1} 页 ---\n{text}")

    if pages:
        return "\n\n".join(pages)

    # 文本为空 → 尝试 OCR 回退
    return _read_pdf_ocr(path)


def _read_pdf_ocr(path: Path) -> str:
    """扫描版 PDF 的 OCR 回退，需要 Tesseract + 中文语言包"""
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        raise ValueError(
            "该 PDF 为扫描版图片，无法提取文字。\n"
            "如需 OCR 支持，请安装：pip install pdf2image pytesseract\n"
            "并安装 Tesseract OCR 引擎（含 chi_sim 中文语言包）。"
        )

    try:
        images = convert_from_path(str(path), dpi=200)
    except Exception as e:
        raise ValueError(
            f"PDF 转图片失败（可能缺少 Poppler）：{e}\n"
            "Windows 请下载 Poppler 并添加到 PATH：https://github.com/oschwartz10612/poppler-windows"
        )

    pages = []
    for i, img in enumerate(images):
        try:
            text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        except Exception:
            text = pytesseract.image_to_string(img)
        if text.strip():
            pages.append(f"--- 第 {i + 1} 页（OCR）---\n{text.strip()}")

    if not pages:
        raise ValueError(f"OCR 未能从 PDF 中识别出任何文字：{path.name}")

    return "\n\n".join(pages)


def read_txt(path: Path) -> str:
    for encoding in ["utf-8", "utf-8-sig", "gbk"]:
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError(f"无法识别文件编码：{path.name}")


def read_markdown(path: Path) -> str:
    return read_txt(path)


def _read_docx_xml_fallback(src_path: Path) -> str:
    """直接解析 word/document.xml 提取文本，完全绕过图片加载，用于含损坏图片的 docx"""
    import zipfile
    from lxml import etree

    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    with zipfile.ZipFile(str(src_path), "r") as z:
        try:
            xml_bytes = z.read("word/document.xml")
        except Exception as e:
            raise ValueError(f"无法读取 docx 文档内容：{src_path.name}（{e}）")

    root = etree.fromstring(xml_bytes)
    paragraphs = []
    for para_elem in root.iter(f"{{{W}}}p"):
        texts = []
        has_mark = False
        for r_elem in para_elem.iter(f"{{{W}}}r"):
            rpr = r_elem.find(f"{{{W}}}rPr")
            if rpr is not None:
                if (rpr.find(f"{{{W}}}b") is not None or
                        rpr.find(f"{{{W}}}highlight") is not None):
                    has_mark = True
            for t in r_elem.findall(f"{{{W}}}t"):
                if t.text:
                    texts.append(t.text)
        text = "".join(texts).strip()
        if text:
            prefix = "【文档格式标注重点】" if has_mark else ""
            paragraphs.append(prefix + text)

    return "\n\n".join(paragraphs)


def read_docx(path: Path) -> str:
    """读取 .docx，同时检测加粗/高亮 run 并插入格式标记"""
    from docx import Document

    try:
        doc = Document(path)
    except Exception:
        # docx 含损坏/缺失图片等问题时，回退到直接解析 XML（绕过图片加载）
        return _read_docx_xml_fallback(path)

    paragraphs = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        # 注意：highlight_color 为 'none' 时 python-docx 会抛 ValueError，需捕获
        def _run_has_mark(run):
            if not run.text.strip():
                return False
            if run.bold:
                return True
            try:
                return run.font.highlight_color is not None
            except ValueError:
                return False

        has_format_mark = any(_run_has_mark(run) for run in para.runs)
        prefix = "【文档格式标注重点】" if has_format_mark else ""
        paragraphs.append(prefix + para.text)

    return "\n\n".join(paragraphs)


def read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def read_excel(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                rows.append(line)
        if rows:
            sheets.append(f"--- 工作表：{sheet_name} ---\n" + "\n".join(rows))
    wb.close()
    if not sheets:
        raise ValueError(f"Excel 文件内容为空：{path.name}")
    return "\n\n".join(sheets)


def read_xls(path: Path) -> str:
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        sheets = []
        for sheet in wb.sheets():
            rows = []
            for rx in range(sheet.nrows):
                cells = [str(sheet.cell_value(rx, cx)).strip() for cx in range(sheet.ncols)]
                line = "\t".join(cells).strip()
                if line:
                    rows.append(line)
            if rows:
                sheets.append(f"--- 工作表：{sheet.name} ---\n" + "\n".join(rows))
        if not sheets:
            raise ValueError(f"XLS 文件内容为空：{path.name}")
        return "\n\n".join(sheets)
    except ImportError:
        raise ValueError("读取 .xls 文件需要安装 xlrd：pip install xlrd")


def read_csv(path: Path) -> str:
    for encoding in ["utf-8", "utf-8-sig", "gbk", "gb2312"]:
        try:
            with open(path, newline="", encoding=encoding) as f:
                reader = csv_module.reader(f)
                rows = []
                for row in reader:
                    line = "\t".join(cell.strip() for cell in row).strip()
                    if line:
                        rows.append(line)
            if rows:
                return "\n".join(rows)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError(f"无法识别 CSV 文件编码：{path.name}")


def read_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ValueError("读取 HTML 文件需要安装 beautifulsoup4：pip install beautifulsoup4 lxml")

    raw = None
    for encoding in ["utf-8", "utf-8-sig", "gbk"]:
        try:
            raw = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        raise ValueError(f"无法识别 HTML 文件编码：{path.name}")

    soup = BeautifulSoup(raw, "lxml")
    # 移除不需要的标签
    for tag in soup(["script", "style", "nav", "header", "footer", "aside", "meta", "link"]):
        tag.decompose()

    paragraphs = []
    for elem in soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li", "td", "th"]):
        text = elem.get_text(separator=" ", strip=True)
        if text:
            paragraphs.append(text)

    if not paragraphs:
        # 回退：提取全部文字
        text = soup.get_text(separator="\n", strip=True)
        return text

    return "\n\n".join(paragraphs)


def read_epub(path: Path) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ValueError("读取 EPUB 文件需要安装：pip install ebooklib beautifulsoup4")

    book = epub.read_epub(str(path))
    chapters = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "lxml")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        if text.strip():
            chapters.append(text.strip())

    if not chapters:
        raise ValueError(f"EPUB 文件内容为空：{path.name}")
    return "\n\n".join(chapters)


def read_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ValueError("读取 RTF 文件需要安装 striprtf：pip install striprtf")

    raw = None
    for encoding in ["utf-8", "utf-8-sig", "gbk", "latin-1"]:
        try:
            raw = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        raise ValueError(f"无法识别 RTF 文件编码：{path.name}")

    text = rtf_to_text(raw)
    if not text.strip():
        raise ValueError(f"RTF 文件内容为空：{path.name}")
    return text


def _convert_legacy_office(path: Path, target_ext: str) -> Path:
    """将 .doc/.ppt 旧格式转换为 .docx/.pptx，返回临时文件路径"""
    tmp_dir = Path(tempfile.mkdtemp())

    # 优先尝试 LibreOffice
    for soffice in ["soffice", "soffice.exe",
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", target_ext, "--outdir", str(tmp_dir), str(path)],
                capture_output=True, timeout=60
            )
            converted = tmp_dir / (path.stem + "." + target_ext)
            if converted.exists():
                return converted
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    # Windows 回退：pywin32 COM
    if sys.platform == "win32":
        try:
            import win32com.client
            app_name = "Word.Application" if target_ext == "docx" else "PowerPoint.Application"
            app = win32com.client.Dispatch(app_name)
            app.Visible = False
            out_path = tmp_dir / (path.stem + "." + target_ext)
            if target_ext == "docx":
                doc = app.Documents.Open(str(path.resolve()))
                doc.SaveAs2(str(out_path.resolve()), FileFormat=16)  # 16 = docx
                doc.Close()
            else:
                prs = app.Presentations.Open(str(path.resolve()))
                prs.SaveAs(str(out_path.resolve()), FileFormat=24)  # 24 = pptx
                prs.Close()
            app.Quit()
            if out_path.exists():
                return out_path
        except Exception:
            pass

    raise ValueError(
        f"无法读取 {path.suffix} 旧格式文件：{path.name}\n"
        "请安装 LibreOffice（推荐）或 Microsoft Office，以支持旧版 .doc/.ppt 格式。\n"
        "LibreOffice 下载：https://www.libreoffice.org/"
    )


def read_doc(path: Path) -> str:
    converted = _convert_legacy_office(path, "docx")
    try:
        return read_docx(converted)
    finally:
        try:
            converted.unlink()
            converted.parent.rmdir()
        except OSError:
            pass


def read_ppt(path: Path) -> str:
    converted = _convert_legacy_office(path, "pptx")
    try:
        return read_pptx(converted)
    finally:
        try:
            converted.unlink()
            converted.parent.rmdir()
        except OSError:
            pass


SUPPORTED_EXTENSIONS = {
    ".pdf":  read_pdf,
    ".txt":  read_txt,
    ".md":   read_markdown,
    ".docx": read_docx,
    ".doc":  read_doc,
    ".pptx": read_pptx,
    ".ppt":  read_ppt,
    ".xlsx": read_excel,
    ".xls":  read_xls,
    ".csv":  read_csv,
    ".html": read_html,
    ".htm":  read_html,
    ".epub": read_epub,
    ".rtf":  read_rtf,
}


def read_file(path: Path) -> str:
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型：{ext}（文件：{path.name}）")
    text = SUPPORTED_EXTENSIONS[ext](path)
    if not text.strip():
        raise ValueError(f"文件内容为空或无法提取文本：{path.name}")
    return text


def collect_files(input_path: Path) -> list:
    if input_path.is_file():
        return [input_path]
    elif input_path.is_dir():
        files = []
        for ext in SUPPORTED_EXTENSIONS:
            files.extend(input_path.glob(f"**/*{ext}"))
        return sorted(files)
    else:
        raise FileNotFoundError(f"路径不存在：{input_path}")


# ─────────────────────────────────────────────
# Section 3: AI 分析
# ─────────────────────────────────────────────

SYSTEM_PROMPT = """你是一个专业的教学内容分析助手，代号"助教Alpha"。你的任务不是创作内容，而是像猎人一样精准地从文本中"抓取"信息。

请严格遵守以下工作流程：

### 任务背景
用户会提供两部分信息：
1.  **【预告区】**：老师在PPT开头列出的重点、难点或考点（可能只有关键词）。
2.  **【正文区】**：后续的长篇文章、讲义或PPT详细内容。

### 核心逻辑
1.  **识别模式**：老师在【正文区】可能只划了关键词的开头，或者根本没划线（因为开头已经预告过了）。
2.  **关联原则**：请根据【预告区】的关键词，在【正文区】中寻找语义最相关、解释最完整的段落。即使正文中没有高亮，只要是对应预告内容的详细解释，都要抓取。

### 执行步骤
1.  **第一步 - 拆解**：读取【预告区】，提取出每一个独立的知识点关键词。
2.  **第二步 - 追踪**：带着这些关键词去扫描【正文区】。
    -   如果遇到关键词只出现了一半（比如只划了前3个字），请自动补全为完整的句子或定义。
    -   如果内容跨页，请尽量把相关的解释片段串联起来。
3.  **第三步 - 提取**：只输出原文中对应的句子或段落，不要自己总结，不要丢掉原文的专业术语。

### 输出格式
请严格按照以下 Markdown 格式输出，不要有任何废话：

####  [在此处填入预告区的第一个关键词]
> [在此处填入正文中找到的最完整的原句或段落]

####  [在此处填入预告区的第二个关键词]
> [在此处填入正文中找到的最完整的原句或段落]
"""

USER_PROMPT_TEMPLATE = """{text}"""


def chunk_text(text: str, max_chars: int = 6000) -> list:
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = []
    current_len = 0

    for para in paragraphs:
        para_len = len(para)
        if current_len + para_len > max_chars and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            current_chunk = [para]
            current_len = para_len
        else:
            current_chunk.append(para)
            current_len += para_len

    if current_chunk:
        chunks.append("\n\n".join(current_chunk))

    return chunks


def call_ai(text_chunk: str, max_retries: int = 3) -> str:
    url = API_BASE_URL.rstrip("/") + "/chat/completions"
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": MODEL,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user",   "content": USER_PROMPT_TEMPLATE.format(text=text_chunk)},
        ],
        "temperature": 0.2,
    }
    last_error = None
    for attempt in range(max_retries):
        try:
            resp = requests.post(url, headers=headers, json=payload, timeout=300)
            resp.raise_for_status()
            return resp.json()["choices"][0]["message"]["content"]
        except Exception as e:
            last_error = e
            if attempt < max_retries - 1:
                print(f"  请求失败（第 {attempt + 1} 次），5 秒后重试：{e}")
                time.sleep(5)
    raise last_error


def parse_ai_response(response_text: str) -> list:
    """解析 AI 返回的 Markdown 格式：#### 关键词\n> 原文片段"""
    points = []
    # 匹配 #### 标题 后跟一行或多行 > 引用内容
    pattern = re.compile(
        r'#{3,5}\s+(.+?)\n((?:>.*\n?)+)',
        re.MULTILINE
    )
    for i, m in enumerate(pattern.finditer(response_text)):
        title = m.group(1).strip()
        # 去掉每行开头的 "> " 并合并
        raw_quote = m.group(2)
        lines = [re.sub(r'^>\s?', '', line) for line in raw_quote.splitlines()]
        fragment = "\n".join(lines).strip()
        if title and fragment:
            points.append({
                "序号": i + 1,
                "标题": title,
                "原文片段": fragment,
                "说明": "",
            })
    if not points:
        print("  警告：AI 响应解析失败，跳过此块。")
    return points


def _split_preview_body(full_text: str) -> tuple[str, str]:
    """
    尝试从全文中自动分离预告区和正文区。
    规则：将前 1/5 的段落（或前 10 段，取较小值）视为预告区，其余为正文区。
    若文档只有一段，则预告区置空，全文作正文区。
    """
    paragraphs = [p for p in full_text.split("\n\n") if p.strip()]
    if len(paragraphs) <= 1:
        return "", full_text

    split_at = min(max(1, len(paragraphs) // 5), 10)
    preview = "\n\n".join(paragraphs[:split_at])
    body    = "\n\n".join(paragraphs[split_at:])
    return preview, body


def extract_knowledge_points(full_text: str, progress_callback=None) -> list:
    preview, body = _split_preview_body(full_text)

    # 构建带结构标签的用户输入
    if preview:
        structured = f"【预告区】\n{preview}\n\n【正文区】\n{body}"
    else:
        structured = f"【正文区】\n{body}"

    # 若正文区过长，分块处理（保持预告区不变）
    body_chunks = chunk_text(body)
    total = len(body_chunks)
    print(f"  正文区已分为 {total} 个块，正在逐块分析...")

    all_points = []
    for i, chunk in enumerate(body_chunks):
        print(f"  正在分析第 {i + 1}/{total} 块...")
        if progress_callback:
            progress_callback(i + 1, total)

        chunk_input = f"【预告区】\n{preview}\n\n【正文区】\n{chunk}" if preview else f"【正文区】\n{chunk}"
        try:
            response = call_ai(chunk_input)
            points = parse_ai_response(response)
            for point in points:
                frag = point.get("原文片段", "")
                point["原文片段"] = frag.replace("【文档格式标注重点】", "").strip()
            offset = len(all_points)
            for j, point in enumerate(points):
                point["序号"] = offset + j + 1
            all_points.extend(points)
            print(f"  第 {i + 1} 块提取到 {len(points)} 个知识点")
        except Exception as e:
            print(f"  警告：第 {i + 1} 块处理失败：{e}")
            continue

    return all_points


# ─────────────────────────────────────────────
# Section 4: 生成输出文档
# ─────────────────────────────────────────────

def _apply_highlight_to_runs(para, fragment: str):
    """在段落的 run 级别精确高亮 fragment，保留其他 run 的所有格式。"""
    from docx.oxml.ns import qn
    from lxml import etree

    full_text = "".join(r.text for r in para.runs)
    idx = full_text.find(fragment)
    if idx == -1:
        return

    frag_end = idx + len(fragment)

    # 计算每个 run 在 full_text 中的字符区间
    run_spans = []
    pos = 0
    for run in para.runs:
        run_spans.append((pos, pos + len(run.text), run))
        pos += len(run.text)

    # 构造新的 (text, highlighted, source_run) 列表
    new_runs = []
    for run_start, run_end, run in run_spans:
        if not run.text:
            continue
        # 完全在 fragment 外
        if run_end <= idx or run_start >= frag_end:
            new_runs.append((run.text, False, run))
            continue
        # 完全在 fragment 内
        if run_start >= idx and run_end <= frag_end:
            new_runs.append((run.text, True, run))
            continue
        # 部分重叠：拆分
        if run_start < idx:
            new_runs.append((run.text[:idx - run_start], False, run))
        frag_s = max(idx, run_start) - run_start
        frag_e = min(frag_end, run_end) - run_start
        new_runs.append((run.text[frag_s:frag_e], True, run))
        if run_end > frag_end:
            new_runs.append((run.text[frag_end - run_start:], False, run))

    # 用 XML 操作重建 runs（保留原有 rPr 格式）
    p_elem = para._p
    # 移除现有 w:r 元素
    for r_elem in list(p_elem.findall(qn("w:r"))):
        p_elem.remove(r_elem)

    for text, highlighted, src_run in new_runs:
        if not text:
            continue
        new_r = copy.deepcopy(src_run._r)
        # 更新文本
        t_elems = new_r.findall(qn("w:t"))
        if t_elems:
            t_elems[0].text = text
            # 保留空格
            if text.startswith(" ") or text.endswith(" "):
                t_elems[0].set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
            # 移除多余的 w:t
            for extra in t_elems[1:]:
                new_r.remove(extra)
        else:
            t = etree.SubElement(new_r, qn("w:t"))
            t.text = text
            if text.startswith(" ") or text.endswith(" "):
                t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")

        if highlighted:
            rPr = new_r.find(qn("w:rPr"))
            if rPr is None:
                rPr = etree.Element(qn("w:rPr"))
                new_r.insert(0, rPr)
            # 移除已有 highlight 避免重复
            for old_hl in rPr.findall(qn("w:highlight")):
                rPr.remove(old_hl)
            hl = etree.SubElement(rPr, qn("w:highlight"))
            hl.set(qn("w:val"), "yellow")
            # 同时加粗
            bold_elems = rPr.findall(qn("w:b"))
            if not bold_elems:
                etree.SubElement(rPr, qn("w:b"))

        p_elem.append(new_r)


def _annotate_original_docx(original_path: Path, fragments: list, output_path: Path):
    """直接在原始 docx 上标注，保留所有原格式。"""
    from docx import Document

    doc = Document(original_path)
    # 过滤掉空片段
    valid_fragments = [f for f in fragments if f]

    for para in doc.paragraphs:
        para_text = para.text
        if not para_text.strip():
            continue
        for fragment in valid_fragments:
            if fragment in para_text:
                _apply_highlight_to_runs(para, fragment)
                break  # 每段只标注第一个匹配到的片段

    # 同时处理表格中的单元格
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    para_text = para.text
                    if not para_text.strip():
                        continue
                    for fragment in valid_fragments:
                        if fragment in para_text:
                            _apply_highlight_to_runs(para, fragment)
                            break

    doc.save(output_path)
    print(f"  已生成标注文档（保留原格式）：{output_path.name}")


def create_annotated_doc(full_text: str, knowledge_points: list, output_path: Path,
                         original_path: Path = None):
    from docx import Document
    from docx.enum.text import WD_COLOR_INDEX

    fragments = [kp.get("原文片段", "").strip() for kp in knowledge_points if kp.get("原文片段", "").strip()]

    # 若原始文件是 .docx，直接在原文档上标注以保留格式
    if original_path is not None and original_path.suffix.lower() == ".docx":
        try:
            _annotate_original_docx(original_path, fragments, output_path)
            return
        except Exception:
            # 原始 docx 含损坏图片等无法直接操作时，回退到纯文本重建模式
            print("  提示：原始文档含损坏内容，将以纯文本模式生成标注文档。")

    # 其他格式：从提取的纯文本重建文档（原有逻辑）
    doc = Document()
    doc.add_heading("原文（带知识点标注）", level=1)

    for line in full_text.split("\n"):
        if not line.strip():
            doc.add_paragraph("")
            continue

        para = doc.add_paragraph()
        remaining = line

        best_fragment = None
        best_idx = len(remaining) + 1

        for fragment in fragments:
            idx = remaining.find(fragment)
            if idx != -1 and idx < best_idx:
                best_idx = idx
                best_fragment = fragment

        if best_fragment is not None:
            if best_idx > 0:
                para.add_run(remaining[:best_idx])
            run = para.add_run(best_fragment)
            run.bold = True
            run.font.highlight_color = WD_COLOR_INDEX.YELLOW
            after = remaining[best_idx + len(best_fragment):]
            if after:
                para.add_run(after)
        else:
            para.add_run(remaining)

    doc.save(output_path)
    print(f"  已生成标注文档：{output_path.name}")


def create_summary_doc(knowledge_points: list, source_filename: str, output_path: Path):
    from docx import Document
    from docx.shared import RGBColor

    doc = Document()
    doc.add_heading("知识点摘要", level=1)

    subtitle = doc.add_paragraph()
    run = subtitle.add_run(f"来源文件：{source_filename}")
    run.italic = True
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    doc.add_paragraph()

    if not knowledge_points:
        doc.add_paragraph("未提取到知识点。")
        doc.save(output_path)
        return

    for kp in knowledge_points:
        seq   = kp.get("序号", "?")
        title = kp.get("标题", "（无标题）")
        frag  = kp.get("原文片段", "").strip()

        doc.add_heading(f"{seq}. {title}", level=2)

        if frag:
            quote = doc.add_paragraph(style="Quote")
            quote.add_run(frag).italic = True

        doc.add_paragraph()

    doc.save(output_path)
    print(f"  已生成摘要文档：{output_path.name}")


# ─────────────────────────────────────────────
# Section 5: 主入口
# ─────────────────────────────────────────────

def main():
    if not API_KEY:
        print("错误：未找到 API_KEY，请将 .env.example 复制为 .env 并填写 API_KEY。")
        sys.exit(1)

    if len(sys.argv) < 2:
        print("用法：python main.py <文件路径或文件夹路径>")
        print()
        ext_list = "、".join(e.lstrip(".").upper() for e in SUPPORTED_EXTENSIONS)
        print(f"支持的文件格式：{ext_list}")
        print("示例：")
        print("  python main.py notes.pdf")
        print("  python main.py ./lecture_slides/")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    OUTPUT_DIR.mkdir(exist_ok=True)

    try:
        files = collect_files(input_path)
    except FileNotFoundError as e:
        print(f"错误：{e}")
        sys.exit(1)

    if not files:
        print("警告：未找到支持的文件。")
        sys.exit(0)

    print(f"共找到 {len(files)} 个文件待处理。\n")

    for file_path in files:
        print(f"正在处理：{file_path.name}")
        print("-" * 50)

        try:
            full_text = read_file(file_path)
        except Exception as e:
            print(f"  跳过：{e}\n")
            continue

        print(f"  文本提取完成，共 {len(full_text)} 个字符")

        knowledge_points = extract_knowledge_points(full_text)

        if not knowledge_points:
            print("  警告：未提取到任何知识点，跳过生成文档。\n")
            continue

        print(f"  共提取到 {len(knowledge_points)} 个知识点，正在生成文档...")

        stem = file_path.stem
        annotated_path = OUTPUT_DIR / f"原文_带标注_{stem}.docx"
        summary_path   = OUTPUT_DIR / f"知识点摘要_{stem}.docx"

        create_annotated_doc(full_text, knowledge_points, annotated_path, original_path=file_path)
        create_summary_doc(knowledge_points, file_path.name, summary_path)

        print(f"  处理完成！\n")

    print(f"全部处理完毕。输出文件位于：{OUTPUT_DIR.resolve()}")


if __name__ == "__main__":
    main()
